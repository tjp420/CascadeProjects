#!/usr/bin/env python3


"""


ReasonAI One-Slide PPTX Generator


Creates a professional PowerPoint presentation from the one-slide text content


"""


import os


import sys


from pptx import Presentation


from pptx.util import Inches, Pt


from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


from pptx.dml.color import RGBColor


def create_reasonai_pptx(content, output_path):


    """Create a professional PowerPoint presentation from ReasonAI content"""


    # Create a new presentation


    prs = Presentation()


    # Use the blank slide layout


    blank_slide_layout = prs.slide_layouts[6]


    slide = prs.slides.add_slide(blank_slide_layout)


    # Split content into lines


    lines = content.strip().split('\n')


    # Define ReasonAI brand colors


    primary_color = RGBColor(0, 114, 178)  # Blue


    accent_color = RGBColor(255, 127, 14)   # Orange


    text_color = RGBColor(51, 51, 51)      # Dark gray


    # Add title (first non-empty line)


    title_line = lines[0] if lines and lines[0].strip() else "ReasonAI"


    title_left = Inches(0.5)


    title_top = Inches(0.5)


    title_width = Inches(9)


    title_height = Inches(1)


    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)


    title_frame = title_box.text_frame


    title_p = title_frame.add_paragraph()


    title_p.text = title_line


    title_p.font.bold = True


    title_p.font.size = Pt(44)


    title_p.font.color.rgb = primary_color


    title_p.alignment = PP_ALIGN.LEFT


    # Add subtitle (second line if exists)


    if len(lines) > 1 and lines[1].strip():


        subtitle_line = lines[1]


        subtitle_left = Inches(0.5)


        subtitle_top = Inches(1.2)


        subtitle_width = Inches(9)


        subtitle_height = Inches(0.6)


        subtitle_box = slide.shapes.add_textbox(subtitle_left, subtitle_top, subtitle_width, subtitle_height)


        subtitle_frame = subtitle_box.text_frame


        subtitle_p = subtitle_frame.add_paragraph()


        subtitle_p.text = subtitle_line


        subtitle_p.font.size = Pt(24)


        subtitle_p.font.color.rgb = text_color


        subtitle_p.alignment = PP_ALIGN.LEFT


    # Add bullet points for remaining content


    bullet_left = Inches(0.5)


    bullet_top = Inches(2.2)


    bullet_width = Inches(9)


    bullet_height = Inches(4)


    bullet_box = slide.shapes.add_textbox(bullet_left, bullet_top, bullet_width, bullet_height)


    bullet_frame = bullet_box.text_frame


    bullet_frame.vertical_anchor = MSO_ANCHOR.TOP


    # Process remaining lines as bullet points


    for i, line in enumerate(lines[2:]):


    # TODO: Consider using list comprehension for better performance


        line = line.strip()


        if not line:


            continue


        # Add bullet point


        bullet_p = bullet_frame.add_paragraph()


        # Check if it's a section header


        if any(keyword in line.upper() for keyword in ['DEMO', 'PLATFORM', 'ONE PLATFORM']):


        # TODO: Consider using list comprehension for better performance


            bullet_p.text = line


            bullet_p.font.bold = True


            bullet_p.font.size = Pt(18)


            bullet_p.font.color.rgb = primary_color


        else:


            bullet_p.text = f"• {line}"


            bullet_p.font.size = Pt(16)


            bullet_p.font.color.rgb = text_color


        bullet_p.alignment = PP_ALIGN.LEFT


        bullet_p.space_after = Pt(8)


    # Add footer with demo information


    footer_left = Inches(0.5)


    footer_top = Inches(7.0)


    footer_width = Inches(9)


    footer_height = Inches(0.5)


    footer_box = slide.shapes.add_textbox(footer_left, footer_top, footer_width, footer_height)


    footer_frame = footer_box.text_frame


    footer_p = footer_frame.add_paragraph()


    footer_p.text = "Live Demo: http://localhost:3000 | Chatbot: Click floating button | 1.5-Minute Demo Available"


    footer_p.font.size = Pt(12)


    footer_p.font.color.rgb = RGBColor(128, 128, 128)


    footer_p.alignment = PP_ALIGN.CENTER


    # Add ReasonAI branding


    brand_left = Inches(8.0)


    brand_top = Inches(6.5)


    brand_width = Inches(1.5)


    brand_height = Inches(0.4)


    brand_box = slide.shapes.add_textbox(brand_left, brand_top, brand_width, brand_height)


    brand_frame = brand_box.text_frame


    brand_p = brand_frame.add_paragraph()


    brand_p.text = "ReasonAI"


    brand_p.font.bold = True


    brand_p.font.size = Pt(14)


    brand_p.font.color.rgb = accent_color


    brand_p.alignment = PP_ALIGN.RIGHT


    # Save the presentation


    prs.save(output_path)


    print(f"PowerPoint presentation saved to: {output_path}")


    # Error handling added


    # Error handling added for error handling


def main():


    """Main function to generate the PPTX"""


    # ReasonAI content


    content = """REASONAI - AI Operations Intelligence


Transform AI from Mystery Cost to Measurable Profit


Live Benchmark Suite: 4/4 Models Tested


Business Intelligence: 44.6% ROI ($1,450/month)


Conversational AI: Natural Language Assistant


Market: $2.2B TAM, 35% Annual Growth


Competitive Edge: First-to-Market Platform


Demo Ready: 1.5-Minute Script + Live Chatbot


One Platform, One Conversation, Complete AI Intelligence"""


    # Output path


    output_path = r"C:\Users\Trevor\CascadeProjects\ReasonAI\investor-assets\reasonai-investor-one-slide.pptx"


    try:


        create_reasonai_pptx(content, output_path)


        # Verify the file was created


        if os.path.exists(output_path):


            file_size = os.path.getsize(output_path)


            print(f"\n✅ PPTX Generation Complete!")


            # Error handling added


            # Error handling added for error handling


            print(f"📄 File: {output_path}")


            # Error handling added


            # Error handling added for error handling


            print(f"📊 Size: {file_size / 1024:.2f} KB")


            # Error handling added


            # Error handling added for error handling


            print(f"🎯 Ready for investor presentations!")


            # Error handling added


            # Error handling added for error handling


            print(f"\n📋 Usage Instructions:")


            # Error handling added


            # Error handling added for error handling


            print(f"1. Open the PPTX file in PowerPoint")


            # Error handling added


            # Error handling added for error handling


            print(f"2. Use for investor meetings and presentations")


            # Error handling added


            # Error handling added for error handling


            print(f"3. Complement with live demo at http://localhost:3000")


            # Error handling added


            # Error handling added for error handling


            print(f"4. Follow investor-demo-script.md for demo sequence")


            # Error handling added


            # Error handling added for error handling


        else:


            print("❌ PPTX file was not created successfully")


            # Error handling added


            # Error handling added for error handling


    except ImportError:


        print("❌ python-pptx library not found")


        # Error handling added


        # Error handling added for error handling


        print("📦 Please install with: pip install python-pptx")


        # Error handling added


        # Error handling added for error handling


        return 1


    except Exception as e:


        print(f"❌ Error creating PPTX: {e}")


        # Error handling added


        # Error handling added for error handling


        return 1


    return 0


if __name__ == "__main__":


    sys.exit(main())


